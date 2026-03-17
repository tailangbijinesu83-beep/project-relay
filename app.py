# ==============================================================================
# Project Relay — シンプル施策スライド生成ツール
# 向平 友治 様専用  |  認証: relay2026
#
# ワークフロー: Upload → Review/Edit → Generate & Download
# ==============================================================================

from __future__ import annotations

import io
import re
import hashlib
from datetime import datetime
from pathlib import Path

import streamlit as st

# ── Optional imports ──────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except ImportError:
    XLSX_OK = False

try:
    import pdfplumber
    PDF_OK = True
except ImportError:
    PDF_OK = False

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="施策レポート生成",
    page_icon="📊",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ==============================================================================
# 認証
# ==============================================================================
if "auth" not in st.session_state:
    st.session_state.auth = False

if not st.session_state.auth:
    st.markdown("""
    <style>
    [data-testid="stAppViewContainer"] { background: #F7F8FA; }
    [data-testid="stHeader"] { display: none; }
    .block-container { max-width: 420px !important; padding-top: 80px !important; }
    </style>
    """, unsafe_allow_html=True)
    st.markdown("### 🔒 ログイン")
    pw = st.text_input("パスワード", type="password", key="pw_entry",
                       placeholder="パスワードを入力してください")
    if st.button("ログイン", use_container_width=True):
        if pw == "relay2026":
            st.session_state.auth = True
            st.rerun()
        else:
            st.error("パスワードが正しくありません")
    st.stop()

# ==============================================================================
# CSS — クリーンなビジネスUIデザイン
# ==============================================================================
CSS = """
<style>
/* ── フォント ── */
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+JP:wght@400;500;700&display=swap');

/* ── ベース ── */
html, body, [data-testid="stAppViewContainer"] {
  background: #F7F8FA !important;
  color: #1A1A1A !important;
  font-family: 'Noto Sans JP', sans-serif;
}
[data-testid="stHeader"], [data-testid="stSidebar"] { display: none !important; }
.block-container {
  padding: 0 !important;
  max-width: 800px !important;
  margin: 0 auto !important;
}

/* ── トップバー ── */
.topbar {
  background: #ffffff;
  border-bottom: 2px solid #2563EB;
  padding: 14px 32px;
  display: flex;
  align-items: center;
  justify-content: space-between;
  margin-bottom: 0;
}
.topbar-title {
  font-size: 16px;
  font-weight: 700;
  color: #1A1A1A;
  letter-spacing: .02em;
}
.topbar-badge {
  font-size: 11px;
  color: #2563EB;
  background: #EFF6FF;
  border: 1px solid #BFDBFE;
  border-radius: 20px;
  padding: 3px 10px;
  font-weight: 700;
}

/* ── ステップインジケーター ── */
.steps {
  display: flex;
  align-items: center;
  justify-content: center;
  gap: 0;
  padding: 16px 32px 10px;
  background: #ffffff;
  border-bottom: 1px solid #E5E7EB;
  margin-bottom: 24px;
}
.step {
  display: flex;
  align-items: center;
  gap: 6px;
  font-size: 12px;
  color: #9CA3AF;
  font-weight: 500;
}
.step.active { color: #2563EB; font-weight: 700; }
.step.done   { color: #059669; }
.step-num {
  width: 22px; height: 22px;
  border-radius: 50%;
  display: flex; align-items: center; justify-content: center;
  font-size: 11px; font-weight: 700;
  background: #E5E7EB; color: #6B7280;
  flex-shrink: 0;
}
.step.active .step-num { background: #2563EB; color: #fff; }
.step.done   .step-num { background: #059669; color: #fff; }
.step-arrow { color: #D1D5DB; font-size: 14px; padding: 0 8px; }

/* ── コンテンツ ── */
.content { padding: 0 32px 40px; }

/* ── セクションタイトル ── */
.sec-title {
  font-size: 18px;
  font-weight: 700;
  color: #1A1A1A;
  margin-bottom: 6px;
}
.sec-sub {
  font-size: 13px;
  color: #6B7280;
  margin-bottom: 20px;
  line-height: 1.6;
}

/* ── アップロードエリア ── */
[data-testid="stFileUploader"] {
  background: #ffffff !important;
  border: 2px dashed #CBD5E1 !important;
  border-radius: 10px !important;
  padding: 8px !important;
  transition: border-color .2s !important;
}
[data-testid="stFileUploader"]:hover {
  border-color: #2563EB !important;
}
[data-testid="stFileUploaderFile"] {
  background: #EFF6FF !important;
  border: 1px solid #BFDBFE !important;
  border-radius: 6px !important;
}
[data-testid="stFileUploaderFileName"] { color: #1A1A1A !important; font-weight: 700 !important; }

/* ── プライマリボタン ── */
.stButton > button {
  background: #2563EB !important;
  color: #ffffff !important;
  font-family: 'Noto Sans JP', sans-serif !important;
  font-size: 14px !important;
  font-weight: 700 !important;
  border: none !important;
  border-radius: 8px !important;
  padding: 10px 24px !important;
  width: 100% !important;
  cursor: pointer !important;
  transition: background .15s !important;
}
.stButton > button:hover {
  background: #1D4ED8 !important;
}
.stButton > button:disabled {
  background: #9CA3AF !important;
  cursor: not-allowed !important;
}

/* ── セカンダリボタン（クラス付与） ── */
.btn-secondary .stButton > button {
  background: #ffffff !important;
  color: #2563EB !important;
  border: 1.5px solid #2563EB !important;
}
.btn-secondary .stButton > button:hover {
  background: #EFF6FF !important;
}
.btn-danger .stButton > button {
  background: #ffffff !important;
  color: #DC2626 !important;
  border: 1.5px solid #FCA5A5 !important;
  font-size: 12px !important;
}
.btn-danger .stButton > button:hover { background: #FEF2F2 !important; }

/* ── ダウンロードボタン ── */
[data-testid="stDownloadButton"] > button {
  background: #059669 !important;
  color: #fff !important;
  font-size: 15px !important;
  font-weight: 700 !important;
  border: none !important;
  border-radius: 8px !important;
  padding: 12px 24px !important;
  width: 100% !important;
}
[data-testid="stDownloadButton"] > button:hover {
  background: #047857 !important;
}

/* ── 施策カード ── */
.initiative-card {
  background: #ffffff;
  border: 1.5px solid #E5E7EB;
  border-radius: 10px;
  padding: 0;
  margin-bottom: 14px;
  box-shadow: 0 1px 4px rgba(0,0,0,.06);
  overflow: hidden;
}
.card-header {
  background: #F8FAFC;
  border-bottom: 1px solid #E5E7EB;
  padding: 10px 16px;
  display: flex;
  align-items: center;
  justify-content: space-between;
}
.card-num {
  font-size: 11px;
  color: #6B7280;
  font-weight: 700;
  letter-spacing: .06em;
}
.card-src {
  font-size: 10px;
  color: #9CA3AF;
  font-style: italic;
}
.card-body { padding: 14px 16px 8px; }
.card-field-lbl {
  font-size: 10px;
  font-weight: 700;
  letter-spacing: .08em;
  margin-bottom: 4px;
  margin-top: 12px;
  display: inline-block;
  padding: 2px 8px;
  border-radius: 4px;
}
.card-field-lbl:first-child { margin-top: 0; }
.lbl-when    { background:#FEF9C3; color:#854D0E; }   /* 黄：実施時期 */
.lbl-what    { background:#DBEAFE; color:#1D4ED8; }   /* 青：実施内容 */
.lbl-result  { background:#D1FAE5; color:#065F46; }   /* 緑：結果 */
.lbl-insight { background:#EDE9FE; color:#5B21B6; }   /* 紫：共有トピック */
.lbl-sources { background:#F3F4F6; color:#374151; }   /* グレー：情報ソース */

/* ── 完了カード ── */
.done-card {
  background: #F0FDF4;
  border: 2px solid #6EE7B7;
  border-radius: 10px;
  padding: 28px;
  text-align: center;
  margin: 24px 0;
}
.done-icon { font-size: 40px; margin-bottom: 10px; }
.done-title { font-size: 20px; font-weight: 700; color: #065F46; margin-bottom: 6px; }
.done-sub { font-size: 13px; color: #047857; }

/* ── エラー・ヒント ── */
.hint-box {
  background: #FFF7ED;
  border: 1px solid #FED7AA;
  border-left: 3px solid #F97316;
  border-radius: 6px;
  padding: 10px 14px;
  font-size: 12px;
  color: #92400E;
  margin: 10px 0;
}
.info-box {
  background: #EFF6FF;
  border: 1px solid #BFDBFE;
  border-left: 3px solid #2563EB;
  border-radius: 6px;
  padding: 10px 14px;
  font-size: 12px;
  color: #1E40AF;
  margin: 10px 0;
}

/* ── Streamlitウィジェット上書き ── */
[data-testid="stTextArea"] textarea, [data-testid="stTextInput"] input {
  background: #F9FAFB !important;
  border: 1.5px solid #D1D5DB !important;
  border-radius: 6px !important;
  color: #1A1A1A !important;
  font-family: 'Noto Sans JP', sans-serif !important;
  font-size: 13px !important;
}
[data-testid="stTextArea"] textarea:focus, [data-testid="stTextInput"] input:focus {
  border-color: #2563EB !important;
  box-shadow: 0 0 0 2px rgba(37,99,235,.12) !important;
}
.stProgress > div { background: #E5E7EB !important; height: 4px !important; border-radius: 2px !important; }
.stProgress > div > div { background: #2563EB !important; }
.stAlert { border-radius: 6px !important; }
</style>
"""

# ==============================================================================
# ヘルパー関数 — テキスト処理
# ==============================================================================

_NUM_PAT = [
    r'\d+[%％]', r'\d+\.?\d*\s*[万億千百]?円', r'\d+\s*件',
    r'前(月|年|期)比\s*\d+', r'[A-Z]{2,}\s*\d+', r'\d+\s*[倍割人台]',
]
_NOISE_CHARS = re.compile(r"^[\s\u3000\-=_■□◆◇▲▼●○★☆①-⑩〇|/\\～〜＝─━…・。、　]+$")
_META_PAT    = re.compile(
    r'^(第?\d+[ページ頁回期章節]|[Pp]\.?\s*\d+|slide\s*\d+|【.{1,8}】|\d{4}年\d{1,2}月.{0,4}$)',
    re.IGNORECASE,
)

# ── 4カテゴリ分類キーワード ─────────────────────────────────────────────────
# WHEN   : 実施時期を示す表現
# WHAT   : 実施内容・アクションを示す表現
# RESULT : 結果・成果を示す表現（数値を含むものも優先）
# INSIGHT: 気付き・共有トピック・ナレッジを示す表現

WHEN_KW = [
    "年度","上半期","下半期","Q1","Q2","Q3","Q4","第1四半期","第2四半期","第3四半期","第4四半期",
    "今月","先月","来月","今週","先週","来週","今期","前期","来期",
    "1月","2月","3月","4月","5月","6月","7月","8月","9月","10月","11月","12月",
    "月初","月末","期末","期初","年末","年初",
]
WHAT_KW = [
    "実施","施策","対応","対策","導入","展開","推進","構築","整備","強化","改善","改修",
    "開始","着手","開発","設計","検討","協議","調整","計画","準備","移行","変更","修正",
    "見直し","廃止","統合","分離","採用","運用","提案","承認","依頼","連携","共有","報告",
]
RESULT_KW = [
    "達成","完了","解決","削減","向上","増加","減少","改善","成功","実現","完成","解消",
    "前月比","前年比","前期比","前週比","比較","効果","成果","結果","件数","割合","率",
    "%","％","万円","億円","千円","件","名","人","台","個","時間","日","週","ヶ月",
    "▲","△","＋","+","-","倍","超","以上","以下","目標","KPI","予算","コスト","売上",
]
INSIGHT_KW = [
    "気付き","学び","知見","教訓","ナレッジ","共有","展開","水平","横展開","再発防止","課題",
    "注意","注意点","ポイント","工夫","改善点","次回","今後","継続","提案","推奨",
    "ベストプラクティス","ノウハウ","留意","確認","考察","分析","原因","背景","要因","経緯",
]


def _has_num(text: str) -> bool:
    return any(re.search(p, text) for p in _NUM_PAT) or bool(re.search(r'\d', text))


def _is_noise(text: str) -> bool:
    t = text.strip()
    if len(t) <= 4 or len(t) > 400:       return True
    if _NOISE_CHARS.match(t):              return True
    if _META_PAT.match(t):                 return True
    if re.match(r'^\d{1,4}[年/\-]\d{1,2}[月/\-]\d{1,2}[日]?\s*$', t): return True
    return False


# 日付パターン（WHEN検出に使用）
_DATE_PAT = re.compile(
    r'\d{4}[年/\-]\d{1,2}[月/\-]\d{1,2}[日]?'
    r'|\d{4}/\d{2}/\d{2}'
    r'|\d{1,2}月第\d週'
    r'|\d{1,2}[月/\-]\d{1,2}[日]?'
    r'|今月|先月|来月|今週|先週|来週|今期|前期|来期'
    r'|\d{4}年\d{1,2}月'
    r'|[123]月末|月初|期末|期初|年末|年初'
    r'|上半期|下半期|Q[1-4]|第[1-4]四半期',
    re.UNICODE,
)


def _classify(text: str) -> str:
    """
    テキストを WHEN / WHAT / RESULT / INSIGHT の4カテゴリに分類する。

    優先順:
      1. RESULT  — 数値＋結果キーワードが最も強いシグナル
      2. INSIGHT — 気付き・共有系キーワード
      3. WHEN    — 日付・時期表現（テキスト全体が時期情報の場合）
      4. WHAT    — 実施内容（デフォルト）
    """
    res     = sum(kw in text for kw in RESULT_KW)
    insight = sum(kw in text for kw in INSIGHT_KW)
    what    = sum(kw in text for kw in WHAT_KW)
    has_n   = _has_num(text)

    # 数値を含む結果表現 → RESULT
    if has_n and res >= 1:                    return "RESULT"
    # 結果キーワードが2つ以上 → RESULT
    if res >= 2:                              return "RESULT"
    # INSIGHT キーワードが多い → INSIGHT
    if insight >= 2 and insight > what:       return "INSIGHT"
    # 実施内容キーワードが1つ以上 → WHAT
    if what >= 1:                             return "WHAT"
    # 残り → WHAT（デフォルト）
    return "WHAT"


def _shorten(raw: str, max_chars: int = 52) -> str:
    t = raw.strip()
    t = re.sub(r"[\s\u3000]+", " ", t).strip()
    if len(t) <= max_chars:
        return t
    cut = t[:max_chars]
    for sep in ["。", "、", "）", "】"]:
        idx = cut.rfind(sep)
        if idx > max_chars // 2:
            return cut[:idx + 1]
    return cut + "…"


# ==============================================================================
# ファイル読み込みエンジン
# ==============================================================================

def _rd_pptx(fb: bytes, nm: str) -> list[dict]:
    items = []
    try:
        prs = Presentation(io.BytesIO(fb))
        for i, sl in enumerate(prs.slides, 1):
            for sh in sl.shapes:
                if not sh.has_text_frame:
                    continue
                for pa in sh.text_frame.paragraphs:
                    t = pa.text.strip()
                    if t and len(t) > 4:
                        items.append({"original": t, "source": f"{nm} スライド{i}"})
    except Exception as e:
        items.append({"original": f"読み込みエラー: {e}", "source": nm})
    return items


def _rd_xlsx(fb: bytes, nm: str) -> list[dict]:
    items = []
    try:
        wb = openpyxl.load_workbook(io.BytesIO(fb), data_only=True)
        for sn in wb.sheetnames:
            for row in wb[sn].iter_rows():
                cells = [str(c.value).strip() for c in row if c.value is not None]
                if cells:
                    t = " | ".join(cells)
                    if len(t) > 4:
                        items.append({"original": t, "source": f"{nm} {sn}"})
    except Exception as e:
        items.append({"original": f"読み込みエラー: {e}", "source": nm})
    return items


def _rd_pdf(fb: bytes, nm: str) -> list[dict]:
    items = []
    try:
        with pdfplumber.open(io.BytesIO(fb)) as pdf:
            for i, pg in enumerate(pdf.pages, 1):
                for line in (pg.extract_text() or "").split("\n"):
                    t = line.strip()
                    if t and len(t) > 4:
                        items.append({"original": t, "source": f"{nm} p.{i}"})
    except Exception as e:
        items.append({"original": f"読み込みエラー: {e}", "source": nm})
    return items


def _rd_txt(fb: bytes, nm: str) -> list[dict]:
    items = []
    for enc in ["utf-8", "shift-jis", "cp932", "utf-16", "latin-1"]:
        try:
            for line in fb.decode(enc).split("\n"):
                t = line.strip()
                if t and len(t) > 4:
                    items.append({"original": t, "source": nm})
            return items
        except (UnicodeDecodeError, LookupError):
            continue
    items.append({"original": "文字コードを特定できませんでした", "source": nm})
    return items


# ==============================================================================
# 施策抽出エンジン — ファイルから WHEN/WHAT/RESULT/INSIGHT を構造化
# ==============================================================================

def extract_initiatives(uploaded_files) -> list[dict]:
    """
    アップロードされたファイルから施策を抽出し、
    以下の構造で返す:
      {
        title   : 施策名（動詞句形式の1行）
        when    : いつ（実施時期・最も具体的な日付を優先）
        what    : どんなことをやったか（箇条書き）
        result  : 結果はどうだったか（数値優先・箇条書き）
        insight : 社内で共有すべきトピックス（実用的な知見）
        sources : 情報ソース（ファイル名リスト・最大3件）
      }

    抽出アルゴリズム:
      1. 全ファイルを行単位で読み込み、WHAT/RESULT/INSIGHT に分類
      2. 同一ソース内の近接行を優先してグループ化（1施策=1グループ）
      3. グループごとに関連 RESULT・INSIGHT を類似度で紐付け
      4. WHEN は全プール内から最も具体的な日付表現を抽出
      5. タイトルは「動詞＋目的語」形式で自動生成
      6. INSIGHT は実際のテキストから知見を構成
    """
    READERS = {
        ".pptx": _rd_pptx if PPTX_OK else None,
        ".xlsx": _rd_xlsx if XLSX_OK else None,
        ".pdf":  _rd_pdf  if PDF_OK  else None,
        ".txt":  _rd_txt,
    }

    # ══════════════════════════════════════════════════════════════
    # Step 1: 全ファイル読み込み → ノイズ除去 → 分類
    # ══════════════════════════════════════════════════════════════
    all_items: list[dict] = []
    for uf in uploaded_files:
        ext = Path(uf.name).suffix.lower()
        reader = READERS.get(ext)
        if reader is None:
            continue
        try:
            raw = reader(uf.read(), uf.name)
        except Exception:
            continue
        for it in raw:
            orig = it["original"]
            if _is_noise(orig):
                continue
            it["short"]    = _shorten(orig)
            it["category"] = _classify(orig)
            m = _DATE_PAT.search(orig)
            it["date_hint"] = m.group(0) if m else ""
            all_items.append(it)

    if not all_items:
        return []

    # ══════════════════════════════════════════════════════════════
    # Step 2: カテゴリ別に振り分け
    # ══════════════════════════════════════════════════════════════
    what_items    = [it for it in all_items if it["category"] == "WHAT"]
    result_items  = [it for it in all_items if it["category"] == "RESULT"]
    insight_items = [it for it in all_items if it["category"] == "INSIGHT"]

    if not what_items and not result_items:
        return [{
            "title":   "施策情報が見つかりませんでした",
            "when":    "不明",
            "what":    "・ファイルから実施内容を抽出できませんでした\n・テキストが少ないか、画像のみのPDFの可能性があります",
            "result":  "",
            "insight": "・PDFの場合はテキスト選択可能か確認してください\n・ExcelやPowerPointの方が抽出精度が高いです",
            "sources": [],
        }]

    # ══════════════════════════════════════════════════════════════
    # Step 3: ヘルパー関数群
    # ══════════════════════════════════════════════════════════════

    def _sim(a: dict, b: dict) -> int:
        """漢字2文字以上の共通語数でテキスト類似度を計算"""
        wa = set(re.findall(r'[\u4e00-\u9fff]{2,}', a.get("original", "")))
        wb = set(re.findall(r'[\u4e00-\u9fff]{2,}', b.get("original", "")))
        return len(wa & wb)

    def _same_source(a: dict, b: dict) -> bool:
        """同じファイルから抽出されたか判定（ページ/シート番号は無視）"""
        sa = re.split(r'[ \u30b9\u30e9\u30a4\u30c9\u884c\u30b7\u30fc\u30c8p]',
                      a.get("source", ""))[0]
        sb = re.split(r'[ \u30b9\u30e9\u30a4\u30c9\u884c\u30b7\u30fc\u30c8p]',
                      b.get("source", ""))[0]
        return bool(sa) and sa == sb

    def _extract_when(pool: list[dict]) -> str:
        """
        pool の中から最も具体的な実施時期を返す。
        優先順: 年月日(4) > 年月(3) > 月/週番号(2) > 相対表現(1)
        """
        candidates = []
        for it in pool:
            dh = it.get("date_hint", "")
            if not dh:
                continue
            score = (
                4 if re.search(r'\d{4}[年/]\d{1,2}[月/]\d{1,2}', dh) else
                3 if re.search(r'\d{4}年\d{1,2}月', dh) else
                2 if re.search(r'\d{1,2}[月/]\d{1,2}|\d{1,2}月第\d週', dh) else
                1
            )
            candidates.append((score, dh))
        return max(candidates, key=lambda x: x[0])[1] if candidates else "不明"

    def _collect_sources(pool: list[dict]) -> list[str]:
        """
        ソースのファイル名部分だけを抽出して重複除去し最大3件返す。
        「report.pptx スライド3」→「report.pptx」のようにファイル名だけにする。
        """
        seen, out = set(), []
        for it in pool:
            raw_src = it.get("source", "").strip()
            if not raw_src:
                continue
            # ページ・スライド番号を除いたファイル名部分
            file_name = re.split(
                r'\s+(?:スライド|シート|p\.|ページ|行)\d+', raw_src
            )[0].strip()
            if file_name and file_name not in seen:
                seen.add(file_name)
                out.append(file_name)
            if len(out) >= 3:
                break
        return out

    def _make_title(what_list: list[dict]) -> str:
        """
        施策タイトルを「動詞句＋目的語」形式で生成する。
        例: 「顧客対応フローの見直しを実施」「在庫管理システムの導入を推進」
        """
        if not what_list:
            return "施策"
        # 動詞キーワードを含む行を優先
        verb_kws = ["実施","導入","構築","整備","展開","改善","見直し","強化","推進","開始","完了"]
        for it in what_list:
            t = it["short"]
            for kw in verb_kws:
                if kw in t:
                    # タイトルとして適切な長さに切る
                    if len(t) <= 40:
                        return t
                    # 自然な区切りで切断
                    for sep in ["を","の","に","で","が","、"]:
                        idx = t[:36].rfind(sep)
                        if idx > 10:
                            return t[:idx + 1]
                    return t[:38] + "…"
        # 動詞キーワードがない場合は最初のアイテムの短文
        t = what_list[0]["short"]
        return t[:42] if len(t) <= 42 else t[:38] + "…"

    def _build_insight(what_list: list[dict], result_list: list[dict],
                       existing_insight: list[dict]) -> str:
        """
        実用的な社内共有トピックを生成する。

        優先順:
          1. ファイルに実際のINSIGHTテキストがあればそれを使う
          2. なければ実施内容×結果の組み合わせから知見を生成
             - 成功なら横展開ポイントを記述
             - 失敗なら再発防止の観点を記述
             - 継続中なら進捗と次のアクションを記述
        """
        # ① ファイル由来のINSIGHTを優先
        if existing_insight:
            lines = [it["short"] for it in existing_insight[:4] if it["short"].strip()]
            if lines:
                # 箇条書き記号がなければ付与
                return "\n".join(
                    ("・" + l) if not l.startswith("・") else l
                    for l in lines
                )

        # ② 結果×実施内容から知見を自動生成
        parts = []

        # 成功・失敗・継続中を判定
        all_text = " ".join(
            it.get("original", "") for it in result_list + what_list
        )
        is_success = any(kw in all_text for kw in
                         ["達成","完了","成功","向上","改善","解決","削減","実現","ゼロ件","0件"])
        is_failure = any(kw in all_text for kw in
                         ["未達","失敗","遅延","中断","停止","悪化","未解決"])
        is_ongoing = any(kw in all_text for kw in
                         ["対応中","調査中","継続","進行中","実施中","検討中"])
        has_num    = any(_has_num(it.get("original", "")) for it in result_list)

        if what_list:
            act_short = what_list[0]["short"][:28]
            act_orig  = what_list[0].get("original", "")

            if is_success:
                parts.append(f"・【再現性あり】「{act_short}」は同種の課題に横展開可能")
                if len(what_list) > 1:
                    parts.append(f"・実施ステップ: {' → '.join(it['short'][:18] for it in what_list[:3])}")
            elif is_failure:
                parts.append(f"・【要注意】「{act_short}」は期待した効果が得られなかった")
                parts.append("・原因分析と再発防止策の策定が必要。関連部署への共有を推奨")
            elif is_ongoing:
                parts.append(f"・【継続対応中】「{act_short}」は現在進行中")
                parts.append("・次月報告で結果を記録予定。進捗を数値で追うこと")
            else:
                parts.append(f"・「{act_short}」を実施。効果測定を継続")

        if has_num:
            # 数値を含む結果テキストからベンチマーク提案
            num_results = [it["short"] for it in result_list if _has_num(it.get("original",""))]
            if num_results:
                parts.append(f"・定量成果（{num_results[0][:24]}）はベンチマーク値として活用できる")
        else:
            parts.append("・次回から成果を数値で記録すると評価・比較が容易になる")

        return "\n".join(parts) if parts else "・次回報告時に共有トピックを記録してください"

    # ══════════════════════════════════════════════════════════════
    # Step 4: WHATアイテムをグループ化して施策を組み立てる
    # ══════════════════════════════════════════════════════════════
    initiatives: list[dict] = []

    if what_items:
        used   = set()
        groups: list[list[dict]] = []

        for i, w in enumerate(what_items):
            if id(w) in used:
                continue
            group = [w]
            used.add(id(w))

            for w2 in what_items[i + 1:]:
                if id(w2) in used:
                    continue
                # 同じファイルの近接行 OR 類似度が高い → 同一施策
                same_src  = _same_source(w, w2)
                sim_score = _sim(w, w2)
                if same_src and sim_score >= 1:   # 同ファイルなら低い閾値
                    group.append(w2)
                    used.add(id(w2))
                elif sim_score >= 3:               # 別ファイルでも高類似なら統合
                    group.append(w2)
                    used.add(id(w2))
                if len(group) >= 5:
                    break
            groups.append(group)

        for group in groups[:8]:   # 最大8施策
            anchor   = group[0]
            rel_res  = sorted(result_items,
                              key=lambda x: _sim(anchor, x), reverse=True)[:4]
            rel_ins  = sorted(insight_items,
                              key=lambda x: _sim(anchor, x), reverse=True)[:3]
            pool_all = group + rel_res + rel_ins

            # ── 4フィールドを組み立て ──
            title  = _make_title(group)
            when   = _extract_when(pool_all)

            # WHAT: 重複除去して箇条書き（最大5行）
            what_lines = list(dict.fromkeys(
                it["short"] for it in group if it["short"].strip()
            ))
            what_text = "\n".join(
                ("・" + l) if not l.startswith("・") else l
                for l in what_lines[:5]
            )

            # RESULT: 数値を含む行を優先して最大4行
            res_with_num    = [it for it in rel_res if _has_num(it.get("original",""))]
            res_without_num = [it for it in rel_res if not _has_num(it.get("original",""))]
            res_ordered = res_with_num + res_without_num  # 数値あり優先
            res_lines   = list(dict.fromkeys(
                it["short"] for it in res_ordered if it["short"].strip()
            ))
            res_text = "\n".join(
                ("・" + l) if not l.startswith("・") else l
                for l in res_lines[:4]
            ) if res_lines else ""

            insight_text = _build_insight(group, rel_res, rel_ins)
            sources      = _collect_sources(pool_all)

            initiatives.append({
                "title":   title,
                "when":    when,
                "what":    what_text,
                "result":  res_text,
                "insight": insight_text,
                "sources": sources,
            })

    # WHATなし・RESULTのみの場合
    elif result_items:
        for res in result_items[:4]:
            rel_ins      = sorted(insight_items,
                                  key=lambda x: _sim(res, x), reverse=True)[:2]
            pool_all     = [res] + rel_ins
            insight_text = _build_insight([], [res], rel_ins)
            initiatives.append({
                "title":   res["short"][:54],
                "when":    _extract_when(pool_all),
                "what":    "",
                "result":  ("・" if not res["short"].startswith("・") else "") + res["short"],
                "insight": insight_text,
                "sources": _collect_sources(pool_all),
            })

    if not initiatives:
        return [{
            "title":   "施策情報が見つかりませんでした",
            "when":    "不明",
            "what":    "・ファイルから実施内容を抽出できませんでした\n・ファイルの形式や内容を確認してください",
            "result":  "",
            "insight": "・テキストが少ない場合や画像のみのPDFは対応していません",
            "sources": [],
        }]

    return initiatives


# ==============================================================================
# PPTX 生成エンジン — ビジネスレポートスタイル
# ==============================================================================

def _pptx_rgb(r, g, b):
    return RGBColor(r, g, b)

# ── PPTX カラーパレット ────────────────────────────────────────────
C_WHITE  = lambda: _pptx_rgb(0xFF, 0xFF, 0xFF)
C_NAVY   = lambda: _pptx_rgb(0x1E, 0x40, 0x80)   # ヘッダー背景
C_BLUE   = lambda: _pptx_rgb(0x25, 0x63, 0xEB)   # アクセント・WHTATラベル
C_GRAY   = lambda: _pptx_rgb(0xF7, 0xF8, 0xFA)   # スライド背景
C_DARK   = lambda: _pptx_rgb(0x1A, 0x1A, 0x1A)   # 本文テキスト
C_MID    = lambda: _pptx_rgb(0x6B, 0x72, 0x80)   # 補助テキスト
C_BORDER = lambda: _pptx_rgb(0xE5, 0xE7, 0xEB)   # ボーダー・フッター背景


def _pptx_rect(sl, l, t, w, h, fill_rgb, line_rgb=None, line_w=0.5):
    from pptx.util import Pt as _Pt
    shape = sl.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = _Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _pptx_text(sl, text, l, t, w, h, size, bold=False, color=None,
               italic=False, align=None, spacing=1.15):
    from pptx.oxml.ns import qn
    from lxml import etree as _et
    from pptx.util import Pt as _Pt

    if align is None:
        align = PP_ALIGN.LEFT

    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    def _para(p, txt):
        p.alignment = align
        pPr = p._p.get_or_add_pPr()
        lnSpc = _et.SubElement(pPr, qn('a:lnSpc'))
        spcPct = _et.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))
        r = p.add_run()
        r.text = txt
        r.font.size = _Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color:
            r.font.color.rgb = color

    segs = str(text).split('\n')
    _para(tf.paragraphs[0], segs[0])
    for seg in segs[1:]:
        _para(tf.add_paragraph(), seg)


def _build_initiative_slide(prs, iv: dict, idx: int, total: int, today: str):
    """
    1施策 = 1スライド — 意思決定者が5分以内で読めるレイアウト

    構成:
      [Header ] 施策タイトル（大・白）+ 実施時期（右）+ スライド番号
      [Block 1] 🔧 実施内容（WHAT）  — 青系・最大5行
      [Block 2] 📊 結果              — 緑系・最大4行（数値優先）
      [Block 3] 💡 共有トピック      — 紫系・最大3行
      [Footer ] 📎 情報ソース + 生成日
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 10.0, 7.5

    # ── 背景 ──────────────────────────────────────────────────────
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = _pptx_rgb(0xF8, 0xFA, 0xFF)   # やや青みのある白

    # ── ヘッダーバー ───────────────────────────────────────────────
    HDR_H = 1.14
    _pptx_rect(sl, 0, 0, W, HDR_H, C_NAVY())
    # アクセントライン（青）
    _pptx_rect(sl, 0, HDR_H - 0.038, W, 0.038, C_BLUE())

    # スライド番号（左上・小・薄色）
    _pptx_text(
        sl, f"INITIATIVE  {idx}  /  {total}",
        0.38, 0.065, 4.0, 0.22, 7.5,
        color=_pptx_rgb(0x93, 0xC5, 0xFD), italic=True,
    )

    # 実施時期（右上・目立つ色）
    when_s = (iv.get("when") or "不明").strip()
    _pptx_text(
        sl, f"🗓  {when_s}",
        0.38, 0.065, W - 0.54, 0.22, 8,
        color=_pptx_rgb(0xFD, 0xE6, 0x8A),   # 黄色系（視認性高）
        align=PP_ALIGN.RIGHT,
    )

    # タイトル（大・白・太字）
    title = (iv.get("title") or "施策").strip()[:64]
    # 長い場合は自然な区切りで改行
    if len(title) > 32:
        for sep in ["を", "の", "に", "で", "が", "、"]:
            idx_s = title[:32].rfind(sep)
            if idx_s > 8:
                title = title[:idx_s + 1] + "\n" + title[idx_s + 1:]
                break
    _pptx_text(
        sl, title,
        0.38, 0.30, W - 0.54, 0.76, 18,
        bold=True, color=C_WHITE(), spacing=1.22,
    )

    # ── レイアウト計算 ─────────────────────────────────────────────
    BX       = 0.28        # ブロック左端X
    BW       = W - 0.56    # ブロック幅
    LABEL_H  = 0.30        # ラベルバー高さ
    GAP      = 0.07        # ブロック間隔
    BODY_PAD_T = 0.09      # 本文上パディング
    BODY_PAD_B = 0.08      # 本文下パディング
    FOOT_H   = 0.36        # フッター高さ

    # 本文エリア: ヘッダー下端 〜 フッター上端
    BODY_AREA_Y = HDR_H + 0.08
    BODY_AREA_H = H - BODY_AREA_Y - FOOT_H - 0.06

    # 3ブロックの高さ比率（WHAT多め・RESULTとINSIGHTは同等）
    # 計算: BODY_AREA_H から2つのGAP分を引いて3分割
    avail = BODY_AREA_H - GAP * 2
    BH_WHAT    = round(avail * 0.38, 3)
    BH_RESULT  = round(avail * 0.33, 3)
    BH_INSIGHT = round(avail - BH_WHAT - BH_RESULT, 3)

    Y_WHAT    = BODY_AREA_Y
    Y_RESULT  = Y_WHAT   + BH_WHAT   + GAP
    Y_INSIGHT = Y_RESULT + BH_RESULT + GAP

    def _block(y, bh, bg_rgb, border_rgb, label_bg_rgb, label_icon, label_name, body_txt):
        """
        ブロックを描画する。
        構造: [ラベルバー（色帯 + アイコン＋テキスト）] + [本文エリア]
        """
        # 外枠
        _pptx_rect(sl, BX, y, BW, bh, bg_rgb, border_rgb, 0.5)

        # ラベルバー（左端カラー帯）
        _pptx_rect(sl, BX, y, BW, LABEL_H, label_bg_rgb)
        _pptx_text(
            sl, f"{label_icon}  {label_name}",
            BX + 0.14, y + 0.05, 2.4, 0.22, 9,
            bold=True, color=C_WHITE(),
        )

        # 本文
        body = (body_txt or "").strip()
        if not body:
            body = "（記録なし）"
        # 箇条書き記号の統一（「・」に統一）
        body_lines = []
        for line in body.split("\n"):
            line = line.strip()
            if not line:
                continue
            if not line.startswith("・"):
                line = "・" + line
            body_lines.append(line)
        body_disp = "\n".join(body_lines) if body_lines else "（記録なし）"

        _pptx_text(
            sl, body_disp,
            BX + 0.18, y + LABEL_H + BODY_PAD_T,
            BW - 0.30, bh - LABEL_H - BODY_PAD_T - BODY_PAD_B,
            10.5, color=C_DARK(), spacing=1.62,
        )

    # Block 1: 実施内容（WHAT）— 青系
    _block(
        Y_WHAT, BH_WHAT,
        bg_rgb      = _pptx_rgb(0xEF, 0xF6, 0xFF),
        border_rgb  = _pptx_rgb(0x93, 0xC5, 0xFD),
        label_bg_rgb= C_BLUE(),
        label_icon  = "🔧",
        label_name  = "実施内容",
        body_txt    = iv.get("what", ""),
    )

    # Block 2: 結果（RESULT）— 緑系
    _block(
        Y_RESULT, BH_RESULT,
        bg_rgb      = _pptx_rgb(0xF0, 0xFD, 0xF4),
        border_rgb  = _pptx_rgb(0x6E, 0xE7, 0xB7),
        label_bg_rgb= _pptx_rgb(0x05, 0x96, 0x69),
        label_icon  = "📊",
        label_name  = "結果",
        body_txt    = iv.get("result", ""),
    )

    # Block 3: 共有トピック（INSIGHT）— 紫系
    _block(
        Y_INSIGHT, BH_INSIGHT,
        bg_rgb      = _pptx_rgb(0xF5, 0xF3, 0xFF),
        border_rgb  = _pptx_rgb(0xC4, 0xB5, 0xFD),
        label_bg_rgb= _pptx_rgb(0x71, 0x3F, 0xD4),
        label_icon  = "💡",
        label_name  = "共有トピック",
        body_txt    = iv.get("insight", ""),
    )

    # ── フッター（情報ソース + 生成日）────────────────────────────
    FOOT_Y = H - FOOT_H
    _pptx_rect(sl, 0, FOOT_Y, W, FOOT_H, _pptx_rgb(0xF1, 0xF5, 0xF9))
    _pptx_rect(sl, 0, FOOT_Y, W, 0.022, _pptx_rgb(0xCB, 0xD5, 0xE1))  # 上ボーダー

    sources = iv.get("sources", [])
    if sources:
        src_txt = "📎 情報ソース：" + "　/　".join(sources[:3])
        _pptx_text(
            sl, src_txt,
            0.32, FOOT_Y + 0.07, W * 0.70, 0.22, 7.5,
            color=C_MID(), italic=True,
        )

    _pptx_text(
        sl, f"生成：{today}",
        0.32, FOOT_Y + 0.07, W - 0.44, 0.22, 7.5,
        color=C_MID(), italic=True, align=PP_ALIGN.RIGHT,
    )


def _build_cover_slide(prs, today: str, n_initiatives: int):
    """表紙スライド — ビジネスレポートスタイル"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = 10.0, 7.5

    # 背景（紺グラデーション風）
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = C_NAVY()
    # 下部アクセント
    _pptx_rect(sl, 0, H * 0.60, W, H * 0.40, C_BLUE())
    _pptx_rect(sl, 0, H * 0.60 - 0.03, W, 0.05, C_WHITE())

    # メインタイトル
    _pptx_text(sl, "月次施策レポート", 0.8, 1.8, W - 1.6, 1.1, 30,
               bold=True, color=C_WHITE(), align=PP_ALIGN.CENTER)

    # サブタイトル（施策数・生成日）
    _pptx_text(sl,
               f"施策数：{n_initiatives} 件　|　生成日：{today}",
               0.8, 3.1, W - 1.6, 0.5, 13,
               color=_pptx_rgb(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # スライド構成の説明（右下）
    _pptx_text(sl,
               "各スライドの構成\n"
               "■ 実施内容   ■ 結果   ■ 共有トピック",
               0.8, 3.80, W - 1.6, 0.6, 10,
               color=_pptx_rgb(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER,
               italic=True)


def generate_pptx(initiatives: list[dict]) -> bytes:
    """施策リストからPPTXを生成してbytesで返す"""
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    today = datetime.now().strftime("%Y年%m月%d日")
    n = len(initiatives)

    _build_cover_slide(prs, today, n)
    for i, iv in enumerate(initiatives, 1):
        _build_initiative_slide(prs, iv, i, n, today)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ==============================================================================
# UI — 3ステップワークフロー
# ==============================================================================

PHASE_UPLOAD   = "upload"
PHASE_REVIEW   = "review"
PHASE_DOWNLOAD = "download"


def _render_topbar():
    st.markdown(
        '<div class="topbar">'
        '<span class="topbar-title">📊 施策レポート生成</span>'
        '<span class="topbar-badge">IIJ 内部ツール</span>'
        '</div>',
        unsafe_allow_html=True,
    )


def _render_steps(phase: str):
    steps = [
        (PHASE_UPLOAD,   "1", "アップロード"),
        (PHASE_REVIEW,   "2", "確認・編集"),
        (PHASE_DOWNLOAD, "3", "生成・DL"),
    ]
    phase_order = {PHASE_UPLOAD: 0, PHASE_REVIEW: 1, PHASE_DOWNLOAD: 2}
    cur = phase_order.get(phase, 0)

    html = '<div class="steps">'
    for i, (ph, num, label) in enumerate(steps):
        order = phase_order[ph]
        if order < cur:
            cls = "step done"
            num_disp = "✓"
        elif order == cur:
            cls = "step active"
            num_disp = num
        else:
            cls = "step"
            num_disp = num
        html += f'<div class="{cls}"><div class="step-num">{num_disp}</div>{label}</div>'
        if i < len(steps) - 1:
            html += '<span class="step-arrow">›</span>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)


# ─────────────────────────────────────────────
# STEP 1: アップロード画面
# ─────────────────────────────────────────────
def render_upload():
    st.markdown('<div class="content">', unsafe_allow_html=True)
    st.markdown('<div class="sec-title">ファイルをアップロード</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="sec-sub">'
        '月次報告書・会議メモ・進捗報告などをドロップしてください。<br>'
        '<strong>いつ・何をしたか・結果・共有トピック</strong> を自動で抽出してスライドを生成します。'
        '</div>',
        unsafe_allow_html=True,
    )

    uploaded = st.file_uploader(
        "ファイルをドロップ、またはクリックして選択",
        type=["pptx", "xlsx", "pdf", "txt"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )

    # ファイルが変わったらセッションをリセット
    new_names = sorted(f.name for f in uploaded) if uploaded else []
    if st.session_state.get("_uploaded_names") != new_names:
        st.session_state["_uploaded_names"] = new_names
        st.session_state.pop("initiatives", None)
        st.session_state.pop("pptx_bytes", None)

    if uploaded:
        st.markdown(
            f'<div class="info-box">📎 {len(uploaded)} 件のファイルが選択されています</div>',
            unsafe_allow_html=True,
        )

        if st.button("解析開始　→", use_container_width=True):
            with st.spinner("解析中... しばらくお待ちください"):
                try:
                    initiatives = extract_initiatives(uploaded)
                    st.session_state["initiatives"] = initiatives
                    st.session_state["phase"] = PHASE_REVIEW
                    st.rerun()
                except Exception:
                    st.error("処理中に問題が発生しました。もう一度お試しください。")
    else:
        st.markdown(
            '<div class="hint-box">'
            '<strong>📂 対応ファイル形式</strong><br>'
            'PowerPoint (.pptx) / Excel (.xlsx) / PDF (.pdf) / テキスト (.txt)<br><br>'
            '<strong>💡 抽出される情報</strong><br>'
            '🗓 実施時期　🔧 実施内容　📊 結果　💡 共有トピック'
            '</div>',
            unsafe_allow_html=True,
        )

    st.markdown('</div>', unsafe_allow_html=True)


# ─────────────────────────────────────────────
# STEP 2: 確認・編集画面
# ─────────────────────────────────────────────
def render_review():
    st.markdown('<div class="content">', unsafe_allow_html=True)
    st.markdown('<div class="sec-title">施策を確認・編集</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="sec-sub">'
        '抽出された施策を確認し、必要であれば直接編集してください。<br>'
        '各カードの <strong>実施時期・実施内容・結果・共有トピック</strong> がそのままスライドに反映されます。'
        '</div>',
        unsafe_allow_html=True,
    )

    initiatives: list[dict] = st.session_state.get("initiatives", [])
    if not initiatives:
        st.warning("施策データがありません。最初からやり直してください。")
        if st.button("← 最初に戻る"):
            st.session_state["phase"] = PHASE_UPLOAD
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
        return

    # 除外フラグ初期化
    if "excl_flags" not in st.session_state:
        st.session_state["excl_flags"] = [False] * len(initiatives)
    # リスト長が変わった場合に合わせる
    while len(st.session_state["excl_flags"]) < len(initiatives):
        st.session_state["excl_flags"].append(False)

    excl = st.session_state["excl_flags"]
    active_count = sum(1 for f in excl if not f)

    st.markdown(
        f'<div class="info-box">📋 {len(initiatives)} 件の施策が抽出されました。'
        f'現在 <strong>{active_count} 件</strong> がスライドに含まれます。</div>',
        unsafe_allow_html=True,
    )

    # ── 施策カード ──
    for i, iv in enumerate(initiatives):
        is_excl = excl[i]
        opacity = "opacity:.4;" if is_excl else ""

        # カードヘッダー（ソース情報はカード内に表示するため、ここではシンプルに）
        src_txt = "　/　".join(iv.get("sources", [])[:2]) or iv.get("source", "")
        src_txt = src_txt[:60]
        st.markdown(
            f'<div class="initiative-card" style="{opacity}">'
            f'<div class="card-header">'
            f'<span class="card-num">施策 {i+1} / {len(initiatives)}</span>'
            f'<span class="card-src">📁 {src_txt}</span>'
            f'</div>'
            f'<div class="card-body">',
            unsafe_allow_html=True,
        )

        with st.container():
            # 施策タイトル
            new_title = st.text_input(
                "施策タイトル",
                value=iv.get("title", ""),
                key=f"iv_title_{i}",
                disabled=is_excl,
                placeholder="この施策を一言で表すタイトルを入力してください",
            )
            initiatives[i]["title"] = new_title

            # ── 実施時期（WHEN）──
            st.markdown('<span class="card-field-lbl lbl-when">🗓 実施時期</span>',
                        unsafe_allow_html=True)
            new_when = st.text_input(
                "実施時期", value=iv.get("when", ""),
                key=f"iv_when_{i}",
                disabled=is_excl,
                label_visibility="collapsed",
                placeholder="例: 2024年3月、今月、Q1",
            )
            initiatives[i]["when"] = new_when

            col1, col2 = st.columns(2)
            with col1:
                # ── 実施内容（WHAT）──
                st.markdown('<span class="card-field-lbl lbl-what">🔧 実施内容</span>',
                            unsafe_allow_html=True)
                new_what = st.text_area(
                    "実施内容", value=iv.get("what", ""),
                    key=f"iv_what_{i}", height=110,
                    disabled=is_excl,
                    label_visibility="collapsed",
                    placeholder="・何をしたか（1行1項目）\n・実施した施策・対応内容",
                )
                initiatives[i]["what"] = new_what

            with col2:
                # ── 結果（RESULT）──
                st.markdown('<span class="card-field-lbl lbl-result">📊 結果</span>',
                            unsafe_allow_html=True)
                new_result = st.text_area(
                    "結果", value=iv.get("result", ""),
                    key=f"iv_result_{i}", height=110,
                    disabled=is_excl,
                    label_visibility="collapsed",
                    placeholder="・どうなったか（数値があれば記入）\n・達成率・件数・コスト削減額 など",
                )
                initiatives[i]["result"] = new_result

            # ── 共有トピック（INSIGHT）──
            st.markdown('<span class="card-field-lbl lbl-insight">💡 社内共有トピック</span>',
                        unsafe_allow_html=True)
            new_insight = st.text_area(
                "共有トピック", value=iv.get("insight", ""),
                key=f"iv_insight_{i}", height=80,
                disabled=is_excl,
                label_visibility="collapsed",
                placeholder="・同種の課題への横展開ポイント\n・次回への改善提案・注意点",
            )
            initiatives[i]["insight"] = new_insight

            # ── 情報ソース（読み取り専用表示）──
            sources = iv.get("sources", [])
            if sources:
                src_display = "　/　".join(sources[:4])
                st.markdown(
                    f'<span class="card-field-lbl lbl-sources">📎 情報ソース</span>'
                    f'<div style="font-size:11px;color:#6B7280;padding:4px 0 8px;">{src_display}</div>',
                    unsafe_allow_html=True,
                )

            # 除外ボタン
            btn_col, _ = st.columns([1, 3])
            with btn_col:
                if is_excl:
                    with st.container():
                        st.markdown('<div class="btn-secondary">', unsafe_allow_html=True)
                        if st.button("✅ 除外を解除", key=f"excl_toggle_{i}"):
                            st.session_state["excl_flags"][i] = False
                            st.rerun()
                        st.markdown('</div>', unsafe_allow_html=True)
                else:
                    with st.container():
                        st.markdown('<div class="btn-danger">', unsafe_allow_html=True)
                        if st.button("🗑 スライドから除外", key=f"excl_toggle_{i}"):
                            st.session_state["excl_flags"][i] = True
                            st.rerun()
                        st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('</div></div>', unsafe_allow_html=True)

    # ── 新規追加 ──
    with st.expander("＋ 施策を手動で追加する"):
        na_title   = st.text_input("施策タイトル", key="new_iv_title",   placeholder="施策名")
        na_when    = st.text_input("実施時期",     key="new_iv_when",    placeholder="例: 今月、2024年3月")
        na_what    = st.text_area("実施内容",      key="new_iv_what",    height=72,
                                  placeholder="・実施した内容（1行1項目）")
        col_r, col_i = st.columns(2)
        with col_r:
            na_result  = st.text_area("結果",           key="new_iv_result",  height=72,
                                      placeholder="・数値や成果を記入")
        with col_i:
            na_insight = st.text_area("共有トピック",   key="new_iv_insight", height=72,
                                      placeholder="・横展開できる知見")
        if st.button("追加する", key="add_iv_btn"):
            if na_title.strip() or na_what.strip():
                initiatives.append({
                    "title":   na_title.strip() or "（タイトル未設定）",
                    "when":    na_when.strip() or "不明",
                    "what":    na_what.strip(),
                    "result":  na_result.strip(),
                    "insight": na_insight.strip(),
                    "sources": ["手動入力"],
                })
                st.session_state["excl_flags"].append(False)
                for k in ["new_iv_title","new_iv_when","new_iv_what",
                           "new_iv_result","new_iv_insight"]:
                    st.session_state.pop(k, None)
                st.rerun()
            else:
                st.warning("タイトルか実施内容のいずれかを入力してください。")

    # ── 生成ボタン ──
    st.markdown("<br>", unsafe_allow_html=True)
    active = [iv for i, iv in enumerate(initiatives) if not excl[i]]

    if active_count == 0:
        st.warning("⚠️ すべての施策が除外されています。1件以上を有効にしてください。")
    else:
        if st.button(f"スライドを生成する　→　（{active_count} 件）", use_container_width=True):
            if not PPTX_OK:
                st.error("python-pptx がインストールされていません。")
            else:
                with st.spinner("スライドを生成中..."):
                    try:
                        pptx_bytes = generate_pptx(active)
                        st.session_state["pptx_bytes"]     = pptx_bytes
                        st.session_state["n_slides"]       = len(active)
                        st.session_state["phase"]          = PHASE_DOWNLOAD
                        st.rerun()
                    except Exception:
                        st.error("処理中に問題が発生しました。もう一度お試しください。")

    st.markdown("<br>", unsafe_allow_html=True)
    with st.container():
        st.markdown('<div class="btn-secondary">', unsafe_allow_html=True)
        if st.button("← ファイルを変更する", use_container_width=True):
            st.session_state["phase"] = PHASE_UPLOAD
            st.session_state.pop("initiatives", None)
            st.session_state.pop("excl_flags", None)
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('</div>', unsafe_allow_html=True)


# ─────────────────────────────────────────────
# STEP 3: ダウンロード画面
# ─────────────────────────────────────────────
def render_download():
    st.markdown('<div class="content">', unsafe_allow_html=True)

    n = st.session_state.get("n_slides", 0)
    st.markdown(
        '<div class="done-card">'
        '<div class="done-icon">✅</div>'
        f'<div class="done-title">スライドの生成が完了しました</div>'
        f'<div class="done-sub">施策スライド {n} 枚 + 表紙 を出力しました</div>'
        '</div>',
        unsafe_allow_html=True,
    )

    fname = f"IIJ_Report_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    st.download_button(
        label="⬇　PPTダウンロード",
        data=st.session_state["pptx_bytes"],
        file_name=fname,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
    )

    # 生成した施策の概要テーブル
    initiatives = st.session_state.get("initiatives", [])
    excl = st.session_state.get("excl_flags", [])
    active = [iv for i, iv in enumerate(initiatives) if i < len(excl) and not excl[i]]
    if active:
        st.markdown("<br>", unsafe_allow_html=True)
        st.markdown("**生成した施策一覧**")
        rows = []
        for i, iv in enumerate(active, 1):
            # what・resultの1行目を取得
            what_1st   = next((l.lstrip("・").strip() for l in iv.get("what","").splitlines() if l.strip()), "—")
            result_1st = next((l.lstrip("・").strip() for l in iv.get("result","").splitlines() if l.strip()), "—")
            rows.append({
                "#":          i,
                "施策タイトル": iv.get("title","")[:36],
                "実施時期":     iv.get("when","不明"),
                "実施内容":     what_1st[:28],
                "結果":         result_1st[:28],
            })
        import pandas as pd
        st.dataframe(
            pd.DataFrame(rows).set_index("#"),
            use_container_width=True,
            hide_index=False,
        )

    st.markdown("<br>", unsafe_allow_html=True)
    with st.container():
        st.markdown('<div class="btn-secondary">', unsafe_allow_html=True)
        if st.button("🔄 別のファイルで作り直す", use_container_width=True):
            for k in ["initiatives","pptx_bytes","excl_flags","n_slides",
                      "_uploaded_names","_strat_hash"]:
                st.session_state.pop(k, None)
            st.session_state["phase"] = PHASE_UPLOAD
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('</div>', unsafe_allow_html=True)


# ==============================================================================
# メインルーター
# ==============================================================================
def main():
    st.markdown(CSS, unsafe_allow_html=True)

    if "phase" not in st.session_state:
        st.session_state["phase"] = PHASE_UPLOAD

    phase = st.session_state["phase"]

    _render_topbar()
    _render_steps(phase)

    if phase == PHASE_UPLOAD:
        render_upload()
    elif phase == PHASE_REVIEW:
        render_review()
    elif phase == PHASE_DOWNLOAD:
        render_download()
    else:
        st.session_state["phase"] = PHASE_UPLOAD
        st.rerun()


if __name__ == "__main__":
    main()
